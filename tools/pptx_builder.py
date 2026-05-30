"""
tools/pptx_builder.py
PPTX Builder — assembles completed slide data into a .pptx file
using python-pptx with proper layouts, fonts, and styling.

Themes supported: corporate_blue, midnight, clean_light, forest, warm
"""

import os
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Theme Definitions ─────────────────────────────────────────────────────────
THEMES = {
    "corporate_blue": {
        "bg_title":    RGBColor(0x1A, 0x1A, 0x2E),
        "bg_content":  RGBColor(0xFF, 0xFF, 0xFF),
        "bg_accent":   RGBColor(0x4A, 0x90, 0xD9),
        "text_title":  RGBColor(0xFF, 0xFF, 0xFF),
        "text_body":   RGBColor(0x2C, 0x2C, 0x2C),
        "text_bullet": RGBColor(0x1A, 0x1A, 0x2E),
        "accent":      RGBColor(0x4A, 0x90, 0xD9),
        "bullet_char": "▸",
    },
    "midnight": {
        "bg_title":    RGBColor(0x0D, 0x0D, 0x1A),
        "bg_content":  RGBColor(0x12, 0x12, 0x24),
        "bg_accent":   RGBColor(0x6C, 0x63, 0xFF),
        "text_title":  RGBColor(0xFF, 0xFF, 0xFF),
        "text_body":   RGBColor(0xE0, 0xE0, 0xFF),
        "text_bullet": RGBColor(0xCC, 0xCC, 0xFF),
        "accent":      RGBColor(0x6C, 0x63, 0xFF),
        "bullet_char": "•",
    },
    "clean_light": {
        "bg_title":    RGBColor(0xF5, 0xF7, 0xFA),
        "bg_content":  RGBColor(0xFF, 0xFF, 0xFF),
        "bg_accent":   RGBColor(0x20, 0x80, 0x60),
        "text_title":  RGBColor(0x1A, 0x1A, 0x2E),
        "text_body":   RGBColor(0x33, 0x33, 0x33),
        "text_bullet": RGBColor(0x22, 0x22, 0x22),
        "accent":      RGBColor(0x20, 0x80, 0x60),
        "bullet_char": "→",
    },
    "forest": {
        "bg_title":    RGBColor(0x1B, 0x3A, 0x2D),
        "bg_content":  RGBColor(0xF4, 0xF9, 0xF4),
        "bg_accent":   RGBColor(0x2D, 0x6A, 0x4F),
        "text_title":  RGBColor(0xFF, 0xFF, 0xFF),
        "text_body":   RGBColor(0x1B, 0x3A, 0x2D),
        "text_bullet": RGBColor(0x1B, 0x3A, 0x2D),
        "accent":      RGBColor(0x52, 0xB7, 0x88),
        "bullet_char": "◆",
    },
    "warm": {
        "bg_title":    RGBColor(0x3D, 0x18, 0x0A),
        "bg_content":  RGBColor(0xFF, 0xFB, 0xF7),
        "bg_accent":   RGBColor(0xE0, 0x76, 0x30),
        "text_title":  RGBColor(0xFF, 0xFF, 0xFF),
        "text_body":   RGBColor(0x2C, 0x1A, 0x10),
        "text_bullet": RGBColor(0x3D, 0x18, 0x0A),
        "accent":      RGBColor(0xE0, 0x76, 0x30),
        "bullet_char": "◉",
    },
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text, left, top, width, height,
                 font_size=24, bold=False, color=RGBColor(0,0,0),
                 align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txBox


def _add_accent_bar(slide, color: RGBColor, top=Inches(0.08)):
    """Thin colored bar at top of slide."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, top, SLIDE_W, Inches(0.07)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _build_title_slide(prs, slide_data: dict, theme: dict):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, theme["bg_title"])
    _add_accent_bar(slide, theme["accent"], top=Inches(6.9))

    # Main title
    _add_textbox(
        slide, slide_data["title"],
        Inches(1), Inches(2.2), Inches(11.3), Inches(1.6),
        font_size=44, bold=True, color=theme["text_title"],
        align=PP_ALIGN.CENTER
    )
    # Subtitle
    subtitle = slide_data.get("subtitle") or ""
    if not subtitle and slide_data.get("bullets"):
        subtitle = slide_data["bullets"][0]
    if subtitle:
        _add_textbox(
            slide, subtitle,
            Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.8),
            font_size=20, bold=False, color=RGBColor(0xA8, 0xC8, 0xF0),
            align=PP_ALIGN.CENTER
        )

    # Date bottom right
    date_str = datetime.now().strftime("%B %Y")
    _add_textbox(
        slide, date_str,
        Inches(10), Inches(6.8), Inches(3), Inches(0.4),
        font_size=11, color=RGBColor(0x88, 0x99, 0xAA),
        align=PP_ALIGN.RIGHT
    )


def _build_content_slide(prs, slide_data: dict, theme: dict):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, theme["bg_content"])
    _add_accent_bar(slide, theme["accent"])

    # Slide title
    _add_textbox(
        slide, slide_data["title"],
        Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.85),
        font_size=28, bold=True, color=theme["text_bullet"]
    )

    # Divider line effect via thin rectangle
    div = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(12.33), Inches(0.03))
    div.fill.solid()
    div.fill.fore_color.rgb = theme["accent"]
    div.line.fill.background()

    # Bullets
    bullets = slide_data.get("bullets") or []
    bullet_char = theme["bullet_char"]
    top = Inches(1.3)
    for b in bullets[:5]:
        _add_textbox(
            slide, f"  {bullet_char}  {b}",
            Inches(0.5), top, Inches(12.3), Inches(0.75),
            font_size=18, color=theme["text_bullet"]
        )
        top += Inches(0.82)

    # Speaker note
    note = slide_data.get("speaker_note", "")
    if note:
        slide.notes_slide.notes_text_frame.text = note


def _build_two_column_slide(prs, slide_data: dict, theme: dict):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, theme["bg_content"])
    _add_accent_bar(slide, theme["accent"])

    _add_textbox(
        slide, slide_data["title"],
        Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.85),
        font_size=28, bold=True, color=theme["text_bullet"]
    )

    div = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(12.33), Inches(0.03))
    div.fill.solid()
    div.fill.fore_color.rgb = theme["accent"]
    div.line.fill.background()

    left_bullets = slide_data.get("left_bullets") or slide_data.get("bullets", [])[:3]
    right_bullets = slide_data.get("right_bullets") or slide_data.get("bullets", [])[3:]
    bullet_char = theme["bullet_char"]

    # Left column
    top = Inches(1.35)
    for b in left_bullets[:4]:
        _add_textbox(slide, f" {bullet_char}  {b}", Inches(0.4), top, Inches(6.0), Inches(0.75),
                     font_size=17, color=theme["text_bullet"])
        top += Inches(0.82)

    # Vertical divider
    vdiv = slide.shapes.add_shape(1, Inches(6.6), Inches(1.2), Inches(0.03), Inches(5.5))
    vdiv.fill.solid()
    vdiv.fill.fore_color.rgb = theme["accent"]
    vdiv.line.fill.background()

    # Right column
    top = Inches(1.35)
    for b in right_bullets[:4]:
        _add_textbox(slide, f" {bullet_char}  {b}", Inches(6.8), top, Inches(6.0), Inches(0.75),
                     font_size=17, color=theme["text_bullet"])
        top += Inches(0.82)

    note = slide_data.get("speaker_note", "")
    if note:
        slide.notes_slide.notes_text_frame.text = note


def _build_quote_slide(prs, slide_data: dict, theme: dict):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, theme["bg_accent"])
    _add_accent_bar(slide, theme["bg_title"])

    # Large quote mark
    _add_textbox(slide, "\u201C", Inches(0.5), Inches(0.5), Inches(2), Inches(1.5),
                 font_size=80, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    quote_text = ""
    if slide_data.get("bullets"):
        quote_text = slide_data["bullets"][0]

    _add_textbox(
        slide, quote_text,
        Inches(1.2), Inches(1.8), Inches(10.9), Inches(3.0),
        font_size=26, bold=False, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER
    )
    _add_textbox(
        slide, slide_data["title"],
        Inches(1), Inches(5.6), Inches(11.3), Inches(0.6),
        font_size=16, color=RGBColor(0xDD, 0xEE, 0xFF),
        align=PP_ALIGN.CENTER
    )

    note = slide_data.get("speaker_note", "")
    if note:
        slide.notes_slide.notes_text_frame.text = note


def _build_closing_slide(prs, slide_data: dict, theme: dict):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, theme["bg_title"])
    _add_accent_bar(slide, theme["accent"], top=Inches(6.9))

    _add_textbox(
        slide, slide_data["title"],
        Inches(1), Inches(2.5), Inches(11.3), Inches(1.2),
        font_size=40, bold=True, color=theme["text_title"],
        align=PP_ALIGN.CENTER
    )

    closing_text = ""
    if slide_data.get("bullets"):
        closing_text = slide_data["bullets"][0]
    if closing_text:
        _add_textbox(
            slide, closing_text,
            Inches(2), Inches(4.0), Inches(9.3), Inches(0.8),
            font_size=18, color=RGBColor(0xA8, 0xC8, 0xF0),
            align=PP_ALIGN.CENTER
        )


# ── Main Builder ──────────────────────────────────────────────────────────────

def build_pptx(
    deck_title: str,
    slides: list,
    theme_name: str = "corporate_blue",
    output_dir: str = "generated",
    version: int = 1,
) -> str:
    """
    Assemble the full .pptx from slides data.
    Returns the output file path.
    """
    theme = THEMES.get(theme_name, THEMES["corporate_blue"])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for slide_data in slides:
        layout_type = slide_data.get("layout", "content")

        if layout_type == "title":
            _build_title_slide(prs, slide_data, theme)
        elif layout_type == "two_column":
            _build_two_column_slide(prs, slide_data, theme)
        elif layout_type == "quote":
            _build_quote_slide(prs, slide_data, theme)
        elif layout_type == "closing":
            _build_closing_slide(prs, slide_data, theme)
        else:
            _build_content_slide(prs, slide_data, theme)

    # Save
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"deck_v{version}_{timestamp}.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)

    return filepath
